# filepath: /langchain-rfp/langchain-rfp/src/rfp.py
import os
import asyncio
import json
import time
from pathlib import Path
from typing import List, Dict, Any, Tuple
import logging

import PyPDF2
from docx import Document
from pptx import Presentation
import pandas as pd
from concurrent.futures import ThreadPoolExecutor
import re
import tiktoken

from langchain_openai import AzureChatOpenAI
from langchain_core.messages import HumanMessage

# Logger setup
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Azure OpenAI Client
from functools import lru_cache

@lru_cache()
def get_azure_client():
    AzureChatOpenAI.model_rebuild()
    return AzureChatOpenAI(
        azure_deployment="gpt-4o-mini",
        azure_endpoint="https://usedoai4xnaoa01.openai.azure.com/",
        model="gpt-4o-mini",
        api_version="2024-12-01-preview",
        api_key="5e6ab4a3e9e14c9db15eb6aad6a99ec9",
        openai_api_type="azure",
        temperature=0.0
    )

# Tokenizer
def get_tokenizer(model_name="gpt-4o-mini"):
    try:
        return tiktoken.encoding_for_model(model_name)
    except KeyError:
        return tiktoken.get_encoding("cl100k_base")

# File Extractor
class FileExtractor:
    def __init__(self):
        self.supported_extensions = {'.pdf', '.docx', '.pptx', '.xlsx'}
        self.tokenizer = get_tokenizer()

    def extract_from_pdf(self, file_path: Path) -> List[Dict[str, Any]]:
        text_blocks = []
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text = page.extract_text()
                    if text and text.strip():
                        text_blocks.append(text.strip())
        except Exception as e:
            logger.error(f"Error extracting from PDF {file_path}: {e}")
        return text_blocks

    def extract_from_docx(self, file_path: Path) -> List[Dict[str, Any]]:
        text_blocks = []
        try:
            doc = Document(file_path)
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    text_blocks.append(text)
        except Exception as e:
            logger.error(f"Error extracting from DOCX {file_path}: {e}")
        return text_blocks

    def extract_from_pptx(self, file_path: Path) -> List[Dict[str, Any]]:
        text_blocks = []
        try:
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_blocks.append(shape.text.strip())
        except Exception as e:
            logger.error(f"Error extracting from PPTX {file_path}: {e}")
        return text_blocks

    def extract_from_xlsx(self, file_path: Path) -> List[Dict[str, Any]]:
        text_blocks = []
        try:
            xl_file = pd.ExcelFile(file_path)
            for sheet_name in xl_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                sheet_text = df.to_string(index=False, na_rep='')
                if sheet_text.strip():
                    text_blocks.append(sheet_text)
        except Exception as e:
            logger.error(f"Error extracting from XLSX {file_path}: {e}")
        return text_blocks

    def extract_from_file(self, file_path: Path) -> List[str]:
        extension = file_path.suffix.lower()
        if extension == '.pdf':
            return self.extract_from_pdf(file_path)
        elif extension == '.docx':
            return self.extract_from_docx(file_path)
        elif extension == '.pptx':
            return self.extract_from_pptx(file_path)
        elif extension == '.xlsx':
            return self.extract_from_xlsx(file_path)
        else:
            logger.warning(f"Unsupported file format: {extension}")
            return []

# RFP Processor
class RFPProcessor:
    def __init__(self):
        self.extractor = FileExtractor()
        self.client = get_azure_client()
        self.semaphore = asyncio.Semaphore(16)  
        self.rfp_fields = [
            'Contracting Entity',
            'Project Value',
            'Geographies involved',
            'Project Duration',
            'Project Start Date',
            'Submission Date',
            'Person-months Effort',
            'Scope of work (summary of Scope of Work)',
            'Evaluation Model (QCBS or LCS or QCS or anything else)',
            'Min. eligibility criteria (detailed)',
            'Evaluation Criteria (detailed)',
            'Is sub-contracting allowed?',
            'Limitation of Liability clause',
            'Type of contract: Fixed price/quoted time and means, framework contract'
        ]
        self.rfp_questions = [
            "Who is the Contracting Entity?",
            "What is the Project Value?",
            "Which Geographies are involved?",
            "What is the Project Duration?",
            "What is the Project Start Date?",
            "What is the Submission Date?",
            "How many Person-months Effort is estimated?",
            "What is the Scope of work (summary)?",
            "What is the Evaluation Model used?",
            "What are the Minimum eligibility criteria?",
            "What are the Evaluation Criteria?",
            "Is sub-contracting allowed?",
            "What is the Limitation of Liability clause?",
            "What is the Type of contract? (Fixed price/quoted time and means, framework contract)"
        ]
        self.conversation_history: List[Tuple[str, str]] = []  # memory

    async def extract_files_parallel(self, files: List[Path]) -> List[str]:
        loop = asyncio.get_event_loop()
        with ThreadPoolExecutor(max_workers=10) as executor:
            results = await asyncio.gather(*[
                loop.run_in_executor(executor, self.extractor.extract_from_file, file_path)
                for file_path in files
            ])
        # Flatten the list of lists
        return [item for sublist in results for item in sublist]

    def prepare_merged_text(self, chunks: List[str]) -> str:
        return "\n\n".join(chunks)

    async def ask_questions_batch(self, merged_text: str) -> Dict[str, str]:
        questions_block = "\n".join(
            [f"{i+1}. {q}" for i, q in enumerate(self.rfp_questions)]
        )
        prompt = f"""
            You are an expert RFP analyzer.

            Given the following RFP content, answer each of the following questions in detail.

            Number each answer corresponding to the question number.

            Example:
            1. Detailed answer to question 1

            2. Detailed answer to question 2

            ...

            RFP Content:
            \"\"\"
            {merged_text}
            \"\"\"

            Questions:
            {questions_block}

            Start answering below:
            """

        async with self.semaphore:
            message = HumanMessage(content=prompt)
            response = await asyncio.get_event_loop().run_in_executor(
                None, self.client.invoke, [message]
            )

            response_text = response.content.strip()

            answers = {}
            try:
                parts = re.split(r'\n(?=\d+\.\s)', response_text)  # Split at "1. ", "2. " etc
                for part in parts:
                    match = re.match(r'(\d+)\.\s(.*)', part, re.DOTALL)
                    if match:
                        q_num, ans_text = match.groups()
                        answers[q_num.strip()] = ans_text.strip()
            except Exception as e:
                logger.error(f"Error parsing LLM response: {e}")

            return answers

    async def process_folder(self, folder_path: str) -> Tuple[Dict[str, Any], str]:
        folder = Path(folder_path)
        files = [f for f in folder.iterdir() if f.is_file() and f.suffix.lower() in self.extractor.supported_extensions]

        if not files:
            logger.error("No supported files found in the folder.")
            return {}, ""

        logger.info(f"Extracting text from {len(files)} files...")
        chunks = await self.extract_files_parallel(files)

        if not chunks:
            logger.error("No text extracted from any files.")
            return {}, ""

        merged_text = self.prepare_merged_text(chunks)

        logger.info("Sending batch questions to LLM...")
        answers = await self.ask_questions_batch(merged_text)

        results = {}
        for i, field in enumerate(self.rfp_fields, 1):
            answer = answers.get(str(i), "Not Found")
            results[field] = {
                "value": answer
            }

        return results, merged_text

    async def ask_custom_question(self, merged_text: str, question: str) -> str:
    
        memory_block = ""
        if self.conversation_history:
            memory_block = "\nPrevious Questions and Answers:\n"
            for idx, (q, a) in enumerate(self.conversation_history[-5:], 1):
                memory_block += f"{idx}. Q: {q}\n   A: {a}\n"

        prompt = f"""
    You are an expert in RFP document analysis.

    Here is the complete RFP document content:
    \"\"\"{merged_text}\"\"\"

    {memory_block}

    User's Question:
    {question}

    Provide a detailed, accurate paragraph answer based on the RFP content and prior answers.
    """

        async with self.semaphore:
            message = HumanMessage(content=prompt)
            response = await asyncio.get_event_loop().run_in_executor(
                None, self.client.invoke, [message]
            )
            answer = response.content.strip()

        # I am saving to my conversation history but this is for the session only
        self.conversation_history.append((question, answer))
        return answer


# Main
async def main():
    processor = RFPProcessor()

    FOLDER_PATH = r"C:\Users\MU518RS\MAy 7 RFPs\10. RISL"

    start_time = time.time()
    results, merged_text = await processor.process_folder(FOLDER_PATH)

    if results:
        print("\nSummary of RFP Fields with Detailed Answers:\n")
        for field, answer in results.items():
            print(f"{field}:\n{answer['value']}\n")

        with open('rfp_analysis_combined_folder_detailed.json', 'w') as f:
            json.dump(results, f, indent=2)
        print("\nResults saved to rfp_analysis_combined_folder_detailed.json")

        end_time = time.time()
        print(f"\nTotal Processing Time: {round(end_time - start_time, 2)} seconds")
        print("\n You can now ask your own questions based on the entire folder.")
        while True:
            user_question = input("\n Enter your custom question (or type 'exit' to quit): ")
            if user_question.lower() == 'exit':
                print("Exiting interactive question mode.")
                break

            answer = await processor.ask_custom_question(merged_text, user_question)
            print(f"\n Answer:\n{answer}\n")

    else:
        print("No valid fields extracted.")


if __name__ == "__main__":
    asyncio.run(main())